import os
import uuid
import asyncio
from urllib.parse import urlparse, urljoin
from bs4 import BeautifulSoup
import requests
from docx import Document as DocxDocument
import PyPDF2
import pandas as pd
from pptx import Presentation

JOBS = {}

def safe_filename(url, ext='txt'):
    name = url.replace("https://", "").replace("http://", "").replace("/", "__")
    return f"{name}.{ext}"

def extract_docx(path):
    doc = DocxDocument(path)
    return '\n'.join([p.text for p in doc.paragraphs])

def extract_pdf(path):
    text = ""
    with open(path, "rb") as f:
        reader = PyPDF2.PdfReader(f)
        for page in reader.pages:
            page_text = page.extract_text()
            if page_text:
                text += page_text
    return text

def extract_csv(path):
    df = pd.read_csv(path)
    return df.to_string()

def extract_xlsx(path):
    df = pd.read_excel(path)
    return df.to_string()

def extract_pptx(path):
    prs = Presentation(path)
    text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return text

async def start_scrape_job(url, use_js, doc_types):
    job_id = str(uuid.uuid4())
    domain = urlparse(url).netloc
    out_dir = f"output_{domain}"
    os.makedirs(out_dir, exist_ok=True)
    JOBS[job_id] = {"status": "running", "progress": 0, "domain": domain}
    asyncio.create_task(scrape_site(url, use_js, doc_types, out_dir, job_id))
    return job_id

def get_job_status(job_id):
    job = JOBS.get(job_id, None)
    if job:
        return job["status"], job["progress"]
    else:
        return "not_found", 0

def list_results(domain):
    out_dir = f"output_{domain}"
    return os.listdir(out_dir) if os.path.exists(out_dir) else []

def get_result_file(domain, filename):
    return os.path.join(f"output_{domain}", filename)

async def scrape_site(start_url, use_js, doc_types, out_dir, job_id):
    visited = set()
    queue = [start_url]
    all_files = []
    progress = 0

    while queue:
        url = queue.pop(0)
        if url in visited or not url.startswith("http"):
            continue
        visited.add(url)
        try:
            resp = requests.get(url, timeout=15)
            resp.raise_for_status()
            html = resp.text
        except Exception as e:
            continue

        soup = BeautifulSoup(html, 'html.parser')
        text = soup.get_text()
        title = soup.title.string if soup.title else ""

        filename = safe_filename(url, 'txt')
        with open(os.path.join(out_dir, filename), "w", encoding="utf-8") as f:
            f.write(f"URL: {url}\nTitle: {title}\n\n{text}")

        all_files.append(filename)

        for a_tag in soup.find_all('a', href=True):
            href = urljoin(url, a_tag['href'])
            if is_document_link(href, doc_types):
                try:
                    ext = href.split('.')[-1].lower()
                    doc_filename = safe_filename(href, ext)
                    doc_path = os.path.join(out_dir, doc_filename)
                    r = requests.get(href, timeout=20)
                    with open(doc_path, "wb") as doc_file:
                        doc_file.write(r.content)
                    extracted = ""
                    if ext == "docx":
                        extracted = extract_docx(doc_path)
                    elif ext == "pdf":
                        extracted = extract_pdf(doc_path)
                    elif ext == "csv":
                        extracted = extract_csv(doc_path)
                    elif ext in ["xlsx", "xls"]:
                        extracted = extract_xlsx(doc_path)
                    elif ext == "pptx":
                        extracted = extract_pptx(doc_path)
                    elif ext == "txt":
                        with open(doc_path, "r", encoding="utf-8", errors="ignore") as tf:
                            extracted = tf.read()
                    if extracted:
                        with open(os.path.join(out_dir, doc_filename+".txt"), "w", encoding="utf-8") as ef:
                            ef.write(extracted)
                    all_files.append(doc_filename)
                except Exception as e:
                    continue
            elif is_internal_link(href, urlparse(start_url).netloc) and href not in visited:
                queue.append(href)

        progress = int(len(visited) / (len(visited)+len(queue)+1) * 100)
        JOBS[job_id]["progress"] = progress

    JOBS[job_id]["status"] = "done"
    JOBS[job_id]["progress"] = 100

def is_document_link(link, doc_types):
    return any(link.lower().endswith(f".{ext}") for ext in doc_types)

def is_internal_link(link, base_domain):
    parsed_link = urlparse(link)
    return (parsed_link.netloc == '' or parsed_link.netloc == base_domain)