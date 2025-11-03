import os
import uuid
import asyncio
import aiohttp
from urllib.parse import urlparse, urljoin
from bs4 import BeautifulSoup
from docx import Document as DocxDocument
import PyPDF2
import pandas as pd
from pptx import Presentation
from concurrent.futures import ThreadPoolExecutor
import hashlib

JOBS = {}
MAX_CONCURRENT = 50
MAX_PAGES = 500
executor = ThreadPoolExecutor(max_workers=20)

def safe_filename(url, ext='txt'):
    hash_part = hashlib.md5(url.encode()).hexdigest()[:8]
    name = url.replace("https://", "").replace("http://", "").replace("/", "__")[:150]
    return f"{name}_{hash_part}.{ext}"

def extract_docx(path):
    doc = DocxDocument(path)
    return '\n'.join([p.text for p in doc.paragraphs])

def extract_pdf(path):
    text = ""
    with open(path, "rb") as f:
        reader = PyPDF2.PdfReader(f)
        for page in reader.pages:
            page_text = page.extract_text()
            if page_text:
                text += page_text
    return text

def extract_csv(path):
    df = pd.read_csv(path)
    return df.to_string()

def extract_xlsx(path):
    df = pd.read_excel(path)
    return df.to_string()

def extract_pptx(path):
    prs = Presentation(path)
    text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return text

async def start_scrape_job(url, use_js, doc_types):
    job_id = str(uuid.uuid4())
    domain = urlparse(url).netloc
    out_dir = f"output_{domain}"
    os.makedirs(out_dir, exist_ok=True)
    JOBS[job_id] = {"status": "running", "progress": 0, "domain": domain}
    asyncio.create_task(scrape_site(url, use_js, doc_types, out_dir, job_id))
    return job_id

def get_job_status(job_id):
    job = JOBS.get(job_id, None)
    if job:
        return job["status"], job["progress"]
    else:
        return "not_found", 0

def list_results(domain):
    out_dir = f"output_{domain}"
    return os.listdir(out_dir) if os.path.exists(out_dir) else []

def get_result_file(domain, filename):
    return os.path.join(f"output_{domain}", filename)

async def fetch_url(session, url):
    try:
        async with session.get(url, timeout=aiohttp.ClientTimeout(total=15), ssl=False) as resp:
            if resp.status == 200:
                return await resp.text()
    except:
        pass
    return None

async def download_file(session, url, path):
    try:
        async with session.get(url, timeout=aiohttp.ClientTimeout(total=30), ssl=False) as resp:
            if resp.status == 200:
                with open(path, 'wb') as f:
                    async for chunk in resp.content.iter_chunked(8192):
                        f.write(chunk)
                return True
    except:
        pass
    return False

async def scrape_site(start_url, use_js, doc_types, out_dir, job_id):
    visited = set()
    queue = [start_url]
    base_domain = urlparse(start_url).netloc
    doc_tasks = []
    
    connector = aiohttp.TCPConnector(limit=MAX_CONCURRENT, limit_per_host=20, force_close=False, enable_cleanup_closed=True)
    timeout = aiohttp.ClientTimeout(total=15, connect=5)
    
    async with aiohttp.ClientSession(connector=connector, timeout=timeout) as session:
        while queue and len(visited) < MAX_PAGES:
            batch_size = min(MAX_CONCURRENT, len(queue))
            batch = queue[:batch_size]
            queue = queue[batch_size:]
            
            tasks = []
            for url in batch:
                if url not in visited and url.startswith("http"):
                    visited.add(url)
                    tasks.append(fetch_url(session, url))
            
            results = await asyncio.gather(*tasks, return_exceptions=True)
            
            new_links = []
            for i, html in enumerate(results):
                if html is None or isinstance(html, Exception):
                    continue
                
                url = batch[i]
                try:
                    soup = BeautifulSoup(html, 'lxml')
                except:
                    soup = BeautifulSoup(html, 'html.parser')
                
                text = soup.get_text(separator='\n', strip=True)
                title = soup.title.string if soup.title else ""
                
                filename = safe_filename(url, 'txt')
                with open(os.path.join(out_dir, filename), "w", encoding="utf-8") as f:
                    f.write(f"URL: {url}\nTitle: {title}\n\n{text}")
                
                for a_tag in soup.find_all('a', href=True):
                    href = urljoin(url, a_tag['href'])
                    
                    if is_document_link(href, doc_types):
                        doc_tasks.append(process_document(session, href, doc_types, out_dir))
                    elif is_internal_link(href, base_domain) and href not in visited:
                        new_links.append(href)
                
                JOBS[job_id]["progress"] = min(95, int(len(visited) / MAX_PAGES * 100))
            
            queue.extend([link for link in new_links if link not in queue])
        
        if doc_tasks:
            await asyncio.gather(*doc_tasks, return_exceptions=True)
    
    JOBS[job_id]["status"] = "done"
    JOBS[job_id]["progress"] = 100

async def process_document(session, url, doc_types, out_dir):
    try:
        ext = url.split('.')[-1].lower().split('?')[0]
        if ext not in doc_types:
            return
        
        doc_filename = safe_filename(url, ext)
        doc_path = os.path.join(out_dir, doc_filename)
        
        if await download_file(session, url, doc_path):
            loop = asyncio.get_event_loop()
            extracted = await loop.run_in_executor(executor, extract_document, doc_path, ext)
            
            if extracted:
                with open(os.path.join(out_dir, doc_filename+".txt"), "w", encoding="utf-8") as ef:
                    ef.write(extracted)
    except:
        pass

def extract_document(path, ext):
    try:
        if ext == "docx":
            return extract_docx(path)
        elif ext == "pdf":
            return extract_pdf(path)
        elif ext == "csv":
            return extract_csv(path)
        elif ext in ["xlsx", "xls"]:
            return extract_xlsx(path)
        elif ext == "pptx":
            return extract_pptx(path)
        elif ext == "txt":
            with open(path, "r", encoding="utf-8", errors="ignore") as tf:
                return tf.read()
    except:
        pass
    return ""

def is_document_link(link, doc_types):
    return any(link.lower().split('?')[0].endswith(f".{ext}") for ext in doc_types)

def is_internal_link(link, base_domain):
    parsed_link = urlparse(link)
    return (parsed_link.netloc == '' or parsed_link.netloc == base_domain)
